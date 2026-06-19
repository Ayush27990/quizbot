import os
import json
import re
import time
import logging
import asyncio
import io
import base64
import random

import PyPDF2
from pptx import Presentation
import httpx
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi

from telegram import InlineKeyboardButton, InlineKeyboardMarkup, Update
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    CallbackQueryHandler,
    MessageHandler,
    ContextTypes,
    filters
)
from groq import Groq

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
MEDICINE_GROUP_ID = os.getenv("MEDICINE_GROUP_ID")
ADMIN_ID = 723919716
INTERVAL = 900
PENDING_FILE = "pending_questions.json"

if not TELEGRAM_TOKEN:
    raise ValueError("TELEGRAM_TOKEN missing")
if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY missing")
if not MEDICINE_GROUP_ID:
    raise ValueError("MEDICINE_GROUP_ID missing")

client = Groq(api_key=GROQ_API_KEY)
used_topics = []


# ─────────────────────────────────────────────
# UTILITY HELPERS
# ─────────────────────────────────────────────

def safe_groq_call(model, messages, temperature=0.3, retries=2):
    delay = 2
    last_err = None
    for attempt in range(retries + 1):
        try:
            return client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=temperature
            )
        except Exception as e:
            last_err = e
            msg = str(e).lower()
            if "rate" in msg or "429" in msg or "503" in msg or "timeout" in msg:
                logger.error(f"Groq transient error (attempt {attempt+1}): {e}")
                time.sleep(delay)
                delay *= 2
                continue
            else:
                logger.error(f"Groq error: {e}")
                raise
    raise last_err


def escape_md(text):
    for ch in ["_", "*", "[", "]", "(", ")", "~", "`", ">",
               "#", "+", "-", "=", "|", "{", "}", ".", "!"]:
        text = text.replace(ch, "\\" + ch)
    return text


def extract_json_list(text):
    raw = text
    try:
        text = re.sub(r"```json|```", "", text).strip()
        match = re.search(r"\[.*\]", text, re.DOTALL)
        if not match:
            logger.error("JSON parse: no array found. Raw: " + raw[:500])
            return []
        candidate = match.group()
        try:
            result = json.loads(candidate)
            if isinstance(result, list):
                return result
            return []
        except Exception:
            fixed = re.sub(r",\s*([\]}])", r"\1", candidate)
            fixed = re.sub(r"[\x00-\x1f]+", " ", fixed)
            try:
                result = json.loads(fixed)
                if isinstance(result, list):
                    return result
            except Exception as e2:
                logger.error(f"JSON parse after fix: {e2} | Raw: {raw[:500]}")
            return []
    except Exception as e:
        logger.error(f"JSON parse error: {e} | Raw: {raw[:500]}")
        return []


# ─────────────────────────────────────────────
# KEY FIX: SMART PRE-PROCESSING OF FORWARDED MCQ
# ─────────────────────────────────────────────

def extract_answer_from_explanation(text):
    """
    Look inside the explanation/answer text for clues about which option is correct.
    Returns a letter like 'A', 'B', 'C', or 'D', or None if not found.
    """
    # Pattern 1: "Correct: B" or "Correct answer: B" or "Answer: B"
    m = re.search(r"(?:correct(?:\s+answer)?|answer)\s*[:\-–]\s*([A-D])", text, re.IGNORECASE)
    if m:
        return m.group(1).upper()

    # Pattern 2: "Answer is B" or "The answer is B"
    m = re.search(r"answer\s+is\s+([A-D])", text, re.IGNORECASE)
    if m:
        return m.group(1).upper()

    # Pattern 3: A checkmark emoji right before the option letter
    # e.g.  "✅ B. Plasma ACTH" or "✅B)"
    m = re.search(r"✅\s*([A-D])[\.\)\s]", text)
    if m:
        return m.group(1).upper()

    # Pattern 4: "B is correct" / "Option B is correct"
    m = re.search(r"(?:option\s+)?([A-D])\s+is\s+(?:the\s+)?correct", text, re.IGNORECASE)
    if m:
        return m.group(1).upper()

    # Pattern 5: standalone bold/uppercase "B)" at line start after stripping
    m = re.search(r"^([A-D])[)\.]?\s+is\b", text, re.IGNORECASE | re.MULTILINE)
    if m:
        return m.group(1).upper()

    return None


def normalize_options(text):
    """
    Convert any option format into standard A) B) C) D) lines.
    Handles: A. / A) / 1. / • / plain lines / ✅-prefixed lines.
    Returns (options_list, marked_correct_letter).
    marked_correct_letter is from ✅ marker if present.
    """
    marked_correct = None

    # Remove spoiler tags
    text = text.replace("||", "")

    lines = text.splitlines()
    option_lines = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Detect checkmark on this line
        has_check = "✅" in line
        line_clean = line.replace("✅", "").strip()

        # Match lettered: A. / A) / a. / a)
        m = re.match(r"^([A-Da-d])\s*[.)\-]\s*(.+)", line_clean)
        if m:
            letter = m.group(1).upper()
            content = m.group(2).strip()
            option_lines.append((letter, content))
            if has_check:
                marked_correct = letter
            continue

        # Match numbered: 1. / 1) 
        m = re.match(r"^([1-4])\s*[.)\-]\s*(.+)", line_clean)
        if m:
            idx = int(m.group(1)) - 1
            if 0 <= idx <= 3:
                letter = chr(65 + idx)
                content = m.group(2).strip()
                option_lines.append((letter, content))
                if has_check:
                    marked_correct = letter
            continue

        # Match bullet: • - ► etc followed by text
        m = re.match(r"^[•▪►➤●○◦\-\*]\s*(.+)", line_clean)
        if m:
            content = m.group(1).strip()
            # Only add if it looks like an option (not a long sentence)
            if len(content) < 120 and len(option_lines) < 4:
                letter = chr(65 + len(option_lines))
                option_lines.append((letter, content))
                if has_check:
                    marked_correct = letter

    return option_lines, marked_correct


def preprocess_forwarded_mcq(raw_text):
    """
    Intelligently parse a forwarded MCQ regardless of format.
    Returns a dict with keys: question, options (A)B)C)D)), answer_letter, explanation
    or None if we can't parse it.
    """
    text = raw_text.replace("||", "")

    # ── Split into sections ──────────────────────────────────────────────
    # Try to find explanation block
    explanation = ""
    exp_match = re.search(
        r"(?:💡\s*)?(?:explanation|answer\s+explanation|rationale)\s*[:\-–]?\s*\n(.+)",
        text, re.IGNORECASE | re.DOTALL
    )
    if exp_match:
        explanation = exp_match.group(1).strip()
        text_before_exp = text[:exp_match.start()].strip()
    else:
        text_before_exp = text

    # ── Extract answer letter ────────────────────────────────────────────
    answer_letter = extract_answer_from_explanation(raw_text)  # search full text

    # ── Parse options from the pre-explanation block ─────────────────────
    option_lines, marked_correct = normalize_options(text_before_exp)

    if marked_correct and not answer_letter:
        answer_letter = marked_correct

    # ── Extract question ─────────────────────────────────────────────────
    # Question = everything before the first option line
    question = ""
    if option_lines:
        first_option_letter = option_lines[0][0]
        # Find where options start in the original text
        # Look for the first line that matches an option pattern
        q_lines = []
        for line in text_before_exp.splitlines():
            stripped = line.strip().replace("✅", "").strip()
            # Stop when we hit an option line
            if re.match(r"^[A-Da-d1-4]\s*[.)\-]\s*.+", stripped):
                break
            if re.match(r"^[•▪►➤●○◦\-\*]\s*.+", stripped) and len(q_lines) > 0:
                break
            clean = line.replace("✅", "").strip()
            if clean:
                q_lines.append(clean)
        question = " ".join(q_lines).strip()
    else:
        question = text_before_exp.strip()

    # Remove hashtags and clean
    question = re.sub(r"#\w+", "", question).strip()
    question = re.sub(r"\s{2,}", " ", question)

    if not question or len(option_lines) < 2:
        return None

    # Pad to 4 options if fewer found (shouldn't happen often)
    while len(option_lines) < 4:
        letter = chr(65 + len(option_lines))
        option_lines.append((letter, "N/A"))

    formatted_options = [f"{l}) {c}" for l, c in option_lines[:4]]

    return {
        "question": question,
        "options_text": "\n".join(formatted_options),
        "options_list": formatted_options,
        "answer_letter": answer_letter,   # may be None
        "explanation": explanation
    }


# ─────────────────────────────────────────────
# PENDING PERSISTENCE
# ─────────────────────────────────────────────

def load_pending():
    try:
        with open(PENDING_FILE, "r") as f:
            data = json.load(f)
            for qid, item in data.items():
                if item.get("image_b64"):
                    item["image_bytes"] = base64.b64decode(item["image_b64"])
                else:
                    item["image_bytes"] = None
            return data
    except Exception:
        return {}


def save_pending(data):
    try:
        saveable = {}
        for qid, item in data.items():
            saveable[qid] = {
                "question": item["question"],
                "source": item["source"],
                "image_b64": base64.b64encode(item["image_bytes"]).decode()
                             if item.get("image_bytes") else None
            }
        with open(PENDING_FILE, "w") as f:
            json.dump(saveable, f)
    except Exception as e:
        logger.error(f"Save pending error: {e}")


pending_questions = load_pending()


# ─────────────────────────────────────────────
# TOPIC LIST
# ─────────────────────────────────────────────

HARRISON_TOPICS = [
    "Approach to the patient with chest pain",
    "Acute coronary syndrome STEMI NSTEMI",
    "Heart failure systolic diastolic management",
    "Atrial fibrillation management anticoagulation",
    "Hypertensive emergency urgency",
    "Aortic stenosis clinical features management",
    "Infective endocarditis diagnosis Duke criteria",
    "Pericarditis and cardiac tamponade",
    "Pulmonary embolism diagnosis management",
    "Deep vein thrombosis anticoagulation",
    "Community acquired pneumonia management",
    "Tuberculosis diagnosis treatment",
    "COPD exacerbation management",
    "Asthma acute severe management",
    "Pleural effusion causes diagnosis",
    "Acute respiratory distress syndrome",
    "Pneumothorax types management",
    "Peptic ulcer disease H pylori",
    "Inflammatory bowel disease Crohn ulcerative colitis",
    "Acute pancreatitis severity management",
    "Liver cirrhosis complications management",
    "Hepatitis B C diagnosis treatment",
    "Acute liver failure causes management",
    "Acute kidney injury causes management",
    "Chronic kidney disease complications",
    "Nephrotic syndrome causes management",
    "Nephritic syndrome glomerulonephritis",
    "Diabetic ketoacidosis management",
    "Hyperosmolar hyperglycemic state",
    "Hypothyroidism hyperthyroidism management",
    "Adrenal insufficiency Addison disease",
    "Cushing syndrome diagnosis",
    "Diabetes mellitus type 1 type 2 complications",
    "Hyponatremia hypernatremia management",
    "Hypokalemia hyperkalemia ECG changes",
    "Hypercalcemia hypocalcemia causes",
    "Metabolic acidosis alkalosis approach",
    "Respiratory acidosis alkalosis approach",
    "Anemia approach iron deficiency",
    "Megaloblastic anemia B12 folate",
    "Hemolytic anemia causes workup",
    "Sickle cell disease complications",
    "Thrombocytopenia causes ITP TTP",
    "Disseminated intravascular coagulation",
    "Leukemia acute chronic types",
    "Lymphoma Hodgkin non Hodgkin",
    "Multiple myeloma diagnosis treatment",
    "Rheumatoid arthritis diagnosis management",
    "Systemic lupus erythematosus criteria",
    "Sepsis septic shock management",
    "Meningitis bacterial viral management",
    "Stroke ischemic hemorrhagic management",
    "Seizures epilepsy management",
    "Guillain Barre syndrome",
    "Myasthenia gravis diagnosis treatment",
    "Parkinson disease management",
    "HIV AIDS opportunistic infections",
    "Malaria diagnosis treatment",
    "Typhoid fever diagnosis treatment",
    "Dengue fever management",
    "Approach to fever of unknown origin",
]


# ─────────────────────────────────────────────
# CONTENT FETCHING
# ─────────────────────────────────────────────

def extract_youtube_id(url):
    patterns = [
        r"youtube\.com/watch\?v=([^&]+)",
        r"youtu\.be/([^?]+)",
        r"youtube\.com/shorts/([^?]+)"
    ]
    for pattern in patterns:
        m = re.search(pattern, url)
        if m:
            return m.group(1)
    return None


async def get_youtube_transcript(video_id):
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        return " ".join([t["text"] for t in transcript_list])[:4000]
    except Exception as e:
        logger.error(f"YouTube transcript error: {e}")
        return None


async def fetch_url_content(url):
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        async with httpx.AsyncClient(timeout=15) as c:
            response = await c.get(url, headers=headers, follow_redirects=True)
            soup = BeautifulSoup(response.text, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
            text = re.sub(r"\n{3,}", "\n\n", text)
            return text[:4000]
    except Exception as e:
        logger.error(f"URL fetch error: {e}")
        return None


# ─────────────────────────────────────────────
# QUESTION GENERATION
# ─────────────────────────────────────────────

async def generate_topic():
    used = ", ".join(used_topics[-20:]) if used_topics else "none"
    prompt = (
        "You are a Harrison Internal Medicine expert.\n\n"
        "Suggest ONE specific high-yield Internal Medicine topic.\n\n"
        f"Already used (avoid repeating): {used}\n\n"
        "Must be from Harrison Principles of Internal Medicine.\n"
        "High yield for NEET PG / USMLE / FMGE.\n\n"
        'Return ONLY JSON: {"topic": "Acute coronary syndrome management"}'
    )
    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.9
        )
        text = response.choices[0].message.content
        m = re.search(r"\{.*\}", text, re.DOTALL)
        if m:
            result = json.loads(m.group())
            topic = result.get("topic", "Internal medicine topic")
            used_topics.append(topic)
            if len(used_topics) > 100:
                used_topics.pop(0)
            return topic
        return random.choice(HARRISON_TOPICS)
    except Exception as e:
        logger.error(f"Topic generation error: {e}")
        return random.choice(HARRISON_TOPICS)


async def generate_questions_from_content(content, count=2):
    def build_prompt(c):
        return (
            "You are a Harrison Internal Medicine expert examiner.\n\n"
            f"Generate EXACTLY {count} high-yield clinical MCQs based on:\n\n{c}\n\n"
            "Rules:\n"
            "- Clinical vignette with patient scenario\n"
            "- 4 options, one definitively correct\n"
            "- Detailed explanation\n"
            "- Explain why each wrong option is incorrect\n"
            "- NEET PG / USMLE standard\n\n"
            "Return ONLY a raw JSON array — no markdown, no backticks, no extra text:\n"
            "[\n"
            "  {\n"
            '    "question": "A 55-year-old patient presents with...",\n'
            '    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],\n'
            '    "answer_index": 0,\n'
            '    "explanation": "Correct: A because... B is wrong because..."\n'
            "  }\n"
            "]"
        )
    try:
        response = safe_groq_call(
            "llama-3.3-70b-versatile",
            [{"role": "user", "content": build_prompt(content)}],
            temperature=0.3
        )
        result = extract_json_list(response.choices[0].message.content)
        if result:
            return result

        logger.info("First attempt failed, retrying with trimmed content...")
        trimmed = re.sub(r"[\x00-\x1f]+", " ", content)
        trimmed = re.sub(r"\s{2,}", " ", trimmed).strip()[:2000]
        response2 = safe_groq_call(
            "llama-3.3-70b-versatile",
            [{"role": "user", "content": build_prompt(trimmed)}],
            temperature=0.3
        )
        return extract_json_list(response2.choices[0].message.content)
    except Exception as e:
        logger.error(f"Question generation error: {e}")
        return []


async def rephrase_forwarded_mcq(raw_text):
    """
    The KEY fix: parse the MCQ first to extract answer from explanation,
    then tell the AI exactly what the correct answer is so it never guesses.
    """
    parsed = preprocess_forwarded_mcq(raw_text)

    if parsed:
        # We have structured data — give the AI maximum context
        answer_hint = ""
        if parsed["answer_letter"]:
            answer_hint = (
                f"\n\nIMPORTANT: The correct answer is option {parsed['answer_letter']}. "
                "Use this to set answer_index correctly."
            )

        explanation_hint = ""
        if parsed["explanation"]:
            explanation_hint = f"\n\nOriginal explanation for reference:\n{parsed['explanation']}"

        prompt = (
            "You are a medical MCQ expert.\n\n"
            "Here is a parsed MCQ:\n\n"
            f"QUESTION:\n{parsed['question']}\n\n"
            f"OPTIONS:\n{parsed['options_text']}"
            + answer_hint
            + explanation_hint + "\n\n"
            "Tasks:\n"
            "1. Slightly rephrase the question stem (same clinical meaning)\n"
            "2. Keep the same 4 options in the same order\n"
            "3. Set answer_index to the 0-based index of the correct option "
            "(0=A, 1=B, 2=C, 3=D)\n"
            "4. Write a detailed explanation — why the correct answer is right "
            "and why each other option is wrong\n\n"
            "Return ONLY a raw JSON array — no markdown, no backticks:\n"
            "[\n"
            "  {\n"
            '    "question": "rephrased question...",\n'
            '    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],\n'
            '    "answer_index": 1,\n'
            '    "explanation": "Correct: B because..."\n'
            "  }\n"
            "]"
        )
    else:
        # Fallback: give the AI the raw text but still ask it to find the answer
        prompt = (
            "You are a medical MCQ expert.\n\n"
            "Here is a forwarded MCQ (format may be non-standard):\n\n"
            + raw_text + "\n\n"
            "Tasks:\n"
            "1. Identify the question, the 4 answer options, and the correct answer\n"
            "   - Look for clues: ✅ emoji, 'Correct:', 'Answer:', explanation text\n"
            "2. Slightly rephrase the question stem\n"
            "3. Standardise options as A) B) C) D)\n"
            "4. Set answer_index (0=A,1=B,2=C,3=D) based on what you found\n"
            "5. Write a detailed explanation\n\n"
            "Return ONLY a raw JSON array — no markdown, no backticks:\n"
            "[\n"
            "  {\n"
            '    "question": "rephrased question...",\n'
            '    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],\n'
            '    "answer_index": 0,\n'
            '    "explanation": "Correct: A because..."\n'
            "  }\n"
            "]"
        )

    try:
        response = safe_groq_call(
            "llama-3.3-70b-versatile",
            [{"role": "user", "content": prompt}],
            temperature=0.3
        )
        raw = response.choices[0].message.content
        logger.info("Rephrase raw response: " + raw[:300])
        return extract_json_list(raw)
    except Exception as e:
        logger.error(f"Rephrase error: {e}")
        return []


async def analyze_image(image_bytes):
    try:
        encoded = base64.b64encode(image_bytes).decode()
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": "data:image/jpeg;base64," + encoded}
                    },
                    {
                        "type": "text",
                        "text": (
                            "You are a medical educator.\n"
                            "Analyze this medical image carefully.\n"
                            "Generate 2 high-yield MCQs based on what you see.\n\n"
                            "Return ONLY a raw JSON array — no markdown, no backticks:\n"
                            "[\n"
                            "  {\n"
                            '    "question": "Based on this image...",\n'
                            '    "options": ["A) ...", "B) ...", "C) ...", "D) ..."],\n'
                            '    "answer_index": 0,\n'
                            '    "explanation": "Correct: A because..."\n'
                            "  }\n"
                            "]"
                        )
                    }
                ]
            }]
        )
        return extract_json_list(response.choices[0].message.content)
    except Exception as e:
        logger.error(f"Image analysis error: {e}")
        return []


# ─────────────────────────────────────────────
# APPROVAL + POSTING
# ─────────────────────────────────────────────

async def send_single_for_approval(bot, question, source, image_bytes=None):
    try:
        qid = str(int(time.time())) + "_" + str(random.randint(1000, 9999))
        pending_questions[qid] = {
            "question": question,
            "source": source,
            "image_bytes": image_bytes
        }
        save_pending(pending_questions)

        correct_option = question["options"][question["answer_index"]]
        text = (
            "📋 MCQ FOR APPROVAL\n\n"
            f"📚 Source: {source}\n\n"
            + question["question"] + "\n\n"
            + "\n".join(question["options"])
            + "\n\n✅ Correct: " + correct_option
            + "\n\n💡 Explanation:\n" + question["explanation"]
        )

        keyboard = InlineKeyboardMarkup([
            [
                InlineKeyboardButton("✅ Approve & Post", callback_data="approve_" + qid),
                InlineKeyboardButton("❌ Reject", callback_data="reject_" + qid)
            ],
            [
                InlineKeyboardButton("🔄 Regenerate", callback_data="regen_" + qid)
            ]
        ])

        if len(text) > 4000:
            text = text[:4000] + "...[truncated]"

        if image_bytes:
            await bot.send_photo(
                chat_id=ADMIN_ID,
                photo=io.BytesIO(image_bytes),
                caption="🖼 Image for this MCQ"
            )

        await bot.send_message(
            chat_id=ADMIN_ID,
            text=text,
            reply_markup=keyboard
        )
        logger.info(f"Sent for approval: {source}")
    except Exception as e:
        logger.error(f"Send for approval error: {e}")


async def post_single_to_group(bot, question, image_bytes=None):
    try:
        if image_bytes:
            await bot.send_photo(
                chat_id=MEDICINE_GROUP_ID,
                photo=io.BytesIO(image_bytes),
                caption="🔍 Study this image carefully before answering!"
            )
            await asyncio.sleep(1)

        text_msg = question["question"] + "\n\n" + "\n".join(question["options"])
        await bot.send_message(chat_id=MEDICINE_GROUP_ID, text=text_msg)
        await asyncio.sleep(1)

        clean_options = []
        for opt in question["options"]:
            if len(opt) > 2 and opt[1] == ")":
                clean_options.append(opt[3:].strip())
            else:
                clean_options.append(opt)

        await bot.send_poll(
            chat_id=MEDICINE_GROUP_ID,
            question=question["question"][:300],
            options=clean_options,
            type="quiz",
            correct_option_id=int(question["answer_index"]),
            is_anonymous=True
        )
        await asyncio.sleep(2)

        explanation_escaped = escape_md(question["explanation"])
        spoiler = "💡 Explanation:\n\n||" + explanation_escaped + "||"
        await bot.send_message(
            chat_id=MEDICINE_GROUP_ID,
            text=spoiler,
            parse_mode="MarkdownV2"
        )
        logger.info("Posted to medicine group")
    except Exception as e:
        logger.error(f"Post to group error: {e}")


# ─────────────────────────────────────────────
# SCHEDULED JOB
# ─────────────────────────────────────────────

async def scheduled_job(context: ContextTypes.DEFAULT_TYPE):
    try:
        logger.info("Running scheduled job...")
        topic = await generate_topic()
        logger.info(f"Topic: {topic}")
        questions = await generate_questions_from_content(topic, count=2)
        if not questions:
            logger.error("Failed to generate questions")
            return
        for q in questions:
            await send_single_for_approval(context.bot, q, "Auto: " + topic)
            await asyncio.sleep(1)
    except Exception as e:
        logger.error(f"Scheduled job error: {e}")


# ─────────────────────────────────────────────
# CALLBACK HANDLER
# ─────────────────────────────────────────────

async def handle_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data

    if data.startswith("approve_"):
        qid = data.replace("approve_", "")
        item = pending_questions.get(qid)
        logger.info(f"Approve clicked. qid={qid}, found={item is not None}")
        if item:
            await post_single_to_group(context.bot, item["question"], item.get("image_bytes"))
            pending_questions.pop(qid, None)
            save_pending(pending_questions)
            await query.edit_message_text("✅ Posted to medicine group!")
        else:
            await query.edit_message_text("❌ Question expired. Use /postnow to generate new ones.")

    elif data.startswith("reject_"):
        qid = data.replace("reject_", "")
        pending_questions.pop(qid, None)
        save_pending(pending_questions)
        await query.edit_message_text("❌ Rejected.")

    elif data.startswith("regen_"):
        qid = data.replace("regen_", "")
        pending_questions.pop(qid, None)
        save_pending(pending_questions)
        await query.edit_message_text("🔄 Regenerating...")
        topic = await generate_topic()
        questions = await generate_questions_from_content(topic, count=1)
        if questions:
            await send_single_for_approval(context.bot, questions[0], "Regenerated: " + topic)
        else:
            await context.bot.send_message(chat_id=ADMIN_ID, text="❌ Failed to regenerate. Try /postnow")


# ─────────────────────────────────────────────
# MESSAGE HANDLERS
# ─────────────────────────────────────────────

async def handle_forwarded_poll(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    try:
        poll = update.message.poll
        if not poll:
            return

        question = re.sub(r"#\w+", "", poll.question).strip()
        options = [re.sub(r"#\w+", "", opt.text).strip() for opt in poll.options]

        # Build text in a format our parser can handle
        text = question + "\n\n"
        for i, opt in enumerate(options):
            text += chr(65 + i) + ") " + opt + "\n"

        # Check if poll has correct_option_id (quiz polls)
        if hasattr(poll, "correct_option_id") and poll.correct_option_id is not None:
            correct_letter = chr(65 + poll.correct_option_id)
            text += f"\nAnswer: {correct_letter}"
            logger.info(f"Poll correct option from Telegram: {correct_letter}")

        logger.info("Poll text built: " + text[:300])
        await update.message.reply_text("💬 Forwarded poll detected! Processing...")

        questions = await rephrase_forwarded_mcq(text)
        if not questions:
            await asyncio.sleep(1)
            questions = await rephrase_forwarded_mcq(text)

        if not questions:
            await update.message.reply_text(
                "❌ Could not process poll.\n\nTry /postnow for a fresh question."
            )
            return

        for q in questions:
            await send_single_for_approval(context.bot, q, "Forwarded Poll")
            await asyncio.sleep(1)
    except Exception as e:
        logger.error(f"Forwarded poll error: {e}")
        await update.message.reply_text("❌ Failed to process poll.")


async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return

    text = update.message.text.strip()

    if text.startswith("/debug"):
        await update.message.reply_text("🔍 Raw text received:\n\n" + repr(text[:500]))
        return

    if text.startswith("http://") or text.startswith("https://"):
        youtube_id = extract_youtube_id(text)

        if youtube_id:
            await update.message.reply_text("🎥 YouTube link! Fetching transcript...")
            transcript = await get_youtube_transcript(youtube_id)
            if not transcript:
                await update.message.reply_text("❌ No transcript. Trying page content...")
                transcript = await fetch_url_content(text)
            if not transcript:
                await update.message.reply_text("❌ Could not extract content.")
                return
            await update.message.reply_text("⏳ Generating MCQs from video...")
            questions = await generate_questions_from_content(transcript, count=2)
            source = "YouTube: " + text[:50]
        else:
            await update.message.reply_text("🔗 Article URL! Fetching content...")
            content = await fetch_url_content(text)
            if not content:
                await update.message.reply_text("❌ Could not fetch content.")
                return
            await update.message.reply_text("⏳ Generating MCQs from article...")
            questions = await generate_questions_from_content(content, count=2)
            source = "Article: " + text[:50]

        if not questions:
            await update.message.reply_text("❌ Failed to generate MCQs.")
            return
        for q in questions:
            await send_single_for_approval(context.bot, q, source)
            await asyncio.sleep(1)

    else:
        await update.message.reply_text("💬 Forwarded MCQ text detected! Processing...")
        logger.info("Raw forwarded text: " + text[:500])

        questions = await rephrase_forwarded_mcq(text)
        if not questions:
            logger.info("MCQ rephrase failed, retrying once...")
            await asyncio.sleep(1)
            questions = await rephrase_forwarded_mcq(text)

        if not questions:
            await update.message.reply_text(
                "❌ Could not process MCQ.\n\n"
                "Tip: Forward the full message including the explanation so the bot "
                "can find the correct answer automatically."
            )
            return

        for q in questions:
            await send_single_for_approval(context.bot, q, "Forwarded MCQ")
            await asyncio.sleep(1)


async def handle_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    try:
        await update.message.reply_text("📄 PDF received. Extracting text...")
        file = await update.message.document.get_file()
        file_bytes = await file.download_as_bytearray()
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(bytes(file_bytes)))
        text = ""
        for page in pdf_reader.pages[:15]:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        if not text.strip():
            await update.message.reply_text("❌ Could not extract text.")
            return
        text = text[:6000]
        await update.message.reply_text("⏳ Generating MCQs from PDF...")
        questions = await generate_questions_from_content(text, count=5)
        if not questions:
            questions = await generate_questions_from_content(text, count=2)
        if not questions:
            await update.message.reply_text("❌ Failed to generate MCQs.")
            return
        for q in questions:
            await send_single_for_approval(context.bot, q, "PDF Upload")
            await asyncio.sleep(1)
    except Exception as e:
        logger.error(f"PDF error: {e}")
        await update.message.reply_text("❌ PDF processing failed.")


def get_slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            t = shape.text_frame.text
            if t and t.strip():
                parts.append(t)
        elif hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    combined = "\n".join(parts)
    combined = combined.replace("\x0b", "\n")
    combined = re.sub(r"[\x00-\x09\x0c-\x1f]+", " ", combined)
    combined = re.sub(r"[ \t]{2,}", " ", combined)
    return combined.strip()


def parse_pptx_quiz_blocks(prs, max_blocks=5):
    slides_text = [get_slide_text(s) for s in prs.slides]
    blocks = []
    i = 0
    while i < len(slides_text) and len(blocks) < max_blocks:
        t = slides_text[i]
        option_lines = re.findall(r"(?m)^[A-D][\.\)]\s*.+", t)
        if len(option_lines) >= 3:
            block = t
            consumed = 1
            if i + 1 < len(slides_text) and re.search(r"answer", slides_text[i + 1], re.I):
                block += "\n\n" + slides_text[i + 1]
                consumed = 2
                if i + 2 < len(slides_text) and re.search(r"explanation", slides_text[i + 2], re.I):
                    block += "\n\n" + slides_text[i + 2]
                    consumed = 3
            blocks.append(block)
            i += consumed
        else:
            i += 1
    return blocks, len(slides_text)


async def handle_pptx(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    try:
        await update.message.reply_text("📊 PPTX received. Extracting text...")
        file = await update.message.document.get_file()
        file_bytes = await file.download_as_bytearray()
        prs = Presentation(io.BytesIO(bytes(file_bytes)))

        blocks, total_slides = parse_pptx_quiz_blocks(prs, max_blocks=5)

        if blocks:
            await update.message.reply_text(
                f"✅ Found {len(blocks)} ready-made MCQ(s) in this PPTX "
                f"(out of {total_slides} slides). Processing..."
            )
            sent_any = False
            for block in blocks:
                questions = await rephrase_forwarded_mcq(block)
                if not questions:
                    await asyncio.sleep(1)
                    questions = await rephrase_forwarded_mcq(block)
                if questions:
                    for q in questions:
                        await send_single_for_approval(context.bot, q, "PPTX Upload")
                        await asyncio.sleep(1)
                    sent_any = True
                else:
                    logger.error("Could not process one MCQ block from PPTX")
                await asyncio.sleep(1)
            if not sent_any:
                await update.message.reply_text("❌ Found MCQs but AI failed to process them.")
            return

        text = "\n\n".join(get_slide_text(s) for s in prs.slides)
        text = re.sub(r"[•▪►➤●○◦]", "-", text)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()

        if not text:
            await update.message.reply_text("❌ Could not extract text (images only?).")
            return

        text = text[:6000]
        await update.message.reply_text("⏳ No ready-made MCQs — generating from slide content...")
        questions = await generate_questions_from_content(text, count=5)
        if not questions:
            questions = await generate_questions_from_content(text, count=2)
        if not questions:
            await update.message.reply_text("❌ Failed to generate MCQs from PPTX.")
            return
        for q in questions:
            await send_single_for_approval(context.bot, q, "PPTX Upload")
            await asyncio.sleep(1)
    except Exception as e:
        logger.error(f"PPTX error: {e}")
        await update.message.reply_text(f"❌ PPTX processing failed: {e}")


async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    try:
        await update.message.reply_text("🖼 Image received! Analyzing...")
        if update.message.photo:
            file = await update.message.photo[-1].get_file()
        elif update.message.document:
            file = await update.message.document.get_file()
        else:
            return
        file_bytes = bytes(await file.download_as_bytearray())
        questions = await analyze_image(file_bytes)
        if not questions:
            await update.message.reply_text("❌ Could not generate MCQs from image.")
            return
        await update.message.reply_text("⏳ Sending for approval...")
        for q in questions:
            await send_single_for_approval(context.bot, q, "Image Upload", image_bytes=file_bytes)
            await asyncio.sleep(1)
    except Exception as e:
        logger.error(f"Image error: {e}")
        await update.message.reply_text("❌ Image processing failed.")


# ─────────────────────────────────────────────
# COMMANDS
# ─────────────────────────────────────────────

async def test_group(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    try:
        await context.bot.send_message(chat_id=MEDICINE_GROUP_ID, text="🧪 Test message from bot!")
        await update.message.reply_text("✅ Success! Bot can post to group.")
    except Exception as e:
        await update.message.reply_text(f"❌ Error: {e}")


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    await update.message.reply_text(
        "✅ QuizMasterBot Running!\n\n"
        "Send me:\n"
        "📝 Forwarded MCQ text (any format!)\n"
        "📊 Forwarded MCQ poll\n"
        "📄 PDF\n"
        "📊 PPTX (PowerPoint slides)\n"
        "🖼 Image\n"
        "🔗 Article URL\n"
        "🎥 YouTube URL\n\n"
        "The bot now finds the correct answer from the explanation automatically — "
        "no specific A/B/C/D format needed!\n\n"
        "Auto: 2 Harrison MCQs every 15 min\n\n"
        "Commands:\n"
        "/postnow - Generate immediately\n"
        "/testgroup - Test group posting\n"
        "/status - Check bot status\n"
        "/debug - Echo raw text for troubleshooting"
    )


async def post_now(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    await update.message.reply_text("⏳ Generating Harrison MCQs... please wait")
    await scheduled_job(context)


async def status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    await update.message.reply_text(
        "✅ Bot is running\n"
        f"📊 Pending approvals: {len(pending_questions)}\n"
        f"📚 Topics used: {len(used_topics)}"
    )


async def error_handler(update, context):
    logger.error(f"Update error: {context.error}")


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────

def main():
    logger.info("Starting QuizMasterBot...")
    app = ApplicationBuilder().token(TELEGRAM_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("postnow", post_now))
    app.add_handler(CommandHandler("status", status))
    app.add_handler(CommandHandler("testgroup", test_group))
    app.add_handler(CallbackQueryHandler(handle_callback))
    app.add_handler(MessageHandler(filters.Document.PDF, handle_pdf))
    app.add_handler(MessageHandler(filters.Document.FileExtension("pptx"), handle_pptx))
    app.add_handler(MessageHandler(filters.PHOTO, handle_image))
    app.add_handler(MessageHandler(filters.POLL, handle_forwarded_poll))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    app.add_error_handler(error_handler)

    app.job_queue.run_repeating(scheduled_job, interval=INTERVAL, first=10)

    logger.info("Bot started!")
    app.run_polling(
        allowed_updates=["message", "poll", "poll_answer", "callback_query"],
        drop_pending_updates=True
    )


if __name__ == "__main__":
    main()
